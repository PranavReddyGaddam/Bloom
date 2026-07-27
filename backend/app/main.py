from fastapi import FastAPI, File, UploadFile, Form, HTTPException, Depends, Query, Header, WebSocket, WebSocketDisconnect
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse, Response, RedirectResponse
import asyncio
import logging
import os
import json
import tempfile
from datetime import datetime, timezone
from urllib.parse import urlencode
from typing import Optional, List
from pydantic import BaseModel
import docx
from pptx import Presentation
from dotenv import load_dotenv

from .ai_service import BloomAI
from .models import SummaryRequest, QuizRequest, QuizResponse, SummaryResponse, FlashcardRequest, FlashcardResponse, AnswerCheckRequest, AnswerCheckResponse, AttemptBreakdownResponse, UserStatsResponse, UserAnalyticsResponse, AttemptRecapResponse, RecentAttempt, Subject, CreateSubjectRequest, TutorStartRequest, TutorStartResponse, TutorAnswerRequest, TutorAnswerResponse, TutorWrapRequest, TutorWrapResponse, DocumentInfo, DocumentContent, DueFlashcard, DueFlashcardsResponse, FlashcardReviewRequest, FlashcardReviewResponse, PretestStartRequest, PretestStartResponse, PretestSubmitRequest, PretestSubmitResponse, DueConceptReviewsResponse, PodcastSegment, PodcastResponse, PodcastInfo, DocumentOriginalMeta, RoleplayStartRequest, RoleplayStartResponse, RoleplayEndRequest, RoleplayResultResponse
from . import extraction_agent
from . import url_ingest
from . import tutor_agent
from . import roleplay_agent
from . import stt_service
from . import pretest_agent
from . import memory_service
from . import db
from . import auth
from . import progress
from . import tts_service
from . import storage_service
from . import pdf_render

# Load environment variables
load_dotenv()

# The rest of the backend degrades silently by design — a failed side effect
# must never fail the student's request. That is right, but it means a
# misconfiguration can be invisible; this logger exists for the failures with
# no user-visible symptom at all.
logger = logging.getLogger(__name__)

app = FastAPI(title="Bloom API", version="1.0.0")

# Add CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Allow all origins for development
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Initialize AI service
ai_service = BloomAI()

@app.on_event("shutdown")
async def shutdown():
    await ai_service.aclose()

@app.get("/")
async def root():
    return {"message": "Bloom API is running!"}

@app.get("/health")
async def health_check():
    return {"status": "healthy", "service": "Bloom API"}

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx"}
MAX_UPLOAD_BYTES = 25 * 1024 * 1024

# MIME type per accepted extension, for storing and serving the original file.
# Keyed off the validated extension rather than the browser-supplied
# UploadFile.content_type, which is client-controlled and frequently wrong.
_CONTENT_TYPES = {
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}

@app.get("/progress/{progress_id}")
async def get_progress(progress_id: str):
    """Current stage of a long-running operation, for the frontend's
    progress UI. The id is client-generated and passed alongside the slow
    request; unknown ids just return a null stage (no auth needed — stages
    are generic strings keyed by an unguessable client-side UUID)."""
    return {"stage": progress.get_stage(progress_id)}

@app.post("/upload-pdf")
async def upload_pdf(
    file: UploadFile = File(...),
    progress_id: Optional[str] = Form(None),
    user_id: str = Depends(auth.get_current_user_id),
):
    """Upload and extract text from a PDF, DOCX, or PPTX file"""
    ext = os.path.splitext(file.filename)[1].lower()
    if ext not in SUPPORTED_EXTENSIONS:
        raise HTTPException(status_code=400, detail="Only PDF, DOCX, and PPTX files are allowed")

    content = await file.read()
    if len(content) > MAX_UPLOAD_BYTES:
        raise HTTPException(status_code=413, detail="File too large (max 25 MB)")

    # Upload safety (ROADMAP 5.3): a tempfile-based path never trusts the
    # client-supplied filename (path traversal) and can't collide between
    # concurrent uploads of the same name. Only the sanitized extension is
    # kept so the extractors can dispatch on it.
    file_path = None
    try:
        fd, file_path = tempfile.mkstemp(suffix=ext)
        with os.fdopen(fd, "wb") as buffer:
            buffer.write(content)

        # Extract text based on file type
        if ext == ".pdf":
            text_content = await extraction_agent.extract_structured(
                file_path, ai_service, progress=progress.reporter(progress_id)
            )
        elif ext == ".docx":
            progress.report(progress_id, "Extracting text")
            # Both are sync and CPU-bound; a large deck would otherwise stall
            # every concurrent request, including the /progress polls that
            # exist to keep the UI responsive.
            text_content = await asyncio.to_thread(extract_text_from_docx, file_path)
        else:
            progress.report(progress_id, "Extracting text")
            text_content = await asyncio.to_thread(extract_text_from_pptx, file_path)

        # Memory layer: store this upload in the user's vector memory and
        # surface prior uploads with substantial overlap. Best-effort — a
        # memory failure must never fail the upload itself.
        progress.report(progress_id, "Comparing against your past uploads")
        similar_documents = []
        document_id = None
        try:
            similar_documents, document_id = await memory_service.remember_upload(
                user_id, file.filename, text_content
            )
        except Exception:
            pass

        # Keep the original file so the student can look at the real document,
        # not just its extracted text (extraction is lossy in ways they can't
        # see: title pages lose all but their first line, figures become prose,
        # anything past MAX_ASSEMBLED_CHARS is dropped).
        #
        # Best-effort and separately guarded from the memory step above: a
        # storage failure must neither fail an upload the student waited
        # through nor discard the similar_documents that step just computed.
        # Logged rather than silently passed because this is the one failure
        # here with no user-visible symptom at all — an un-widened S3 policy
        # would leave every document with a null source_key and no error
        # anywhere, which is how the feature ships "working" and does nothing.
        if document_id:
            try:
                key = await storage_service.put_bytes(
                    storage_service.document_key(user_id, document_id, ext),
                    content,
                    _CONTENT_TYPES[ext],
                )
                db.attach_document_source(document_id, key, _CONTENT_TYPES[ext])
            except Exception:
                logger.warning(
                    "Could not store original file for document %s", document_id,
                    exc_info=True,
                )

        return {
            "filename": file.filename,
            "text_content": text_content,
            "word_count": len(text_content.split()),
            "similar_documents": similar_documents,
            "document_id": document_id,
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing file: {str(e)}")
    finally:
        progress.clear(progress_id)
        if file_path is not None:
            try:
                os.remove(file_path)
            except OSError:
                pass

class IngestUrlRequest(BaseModel):
    url: str
    progress_id: Optional[str] = None


@app.post("/ingest-url")
async def ingest_url(
    request: IngestUrlRequest,
    user_id: str = Depends(auth.get_current_user_id),
):
    """Ingest a YouTube video, direct media link, or article URL.

    Returns the same shape as /upload-pdf so the frontend reuses one path:
    a link becomes a document like any other, studiable identically.
    """
    progress_id = request.progress_id
    try:
        result = await url_ingest.ingest_url(
            request.url, ai_service, progress=progress.reporter(progress_id)
        )

        progress.report(progress_id, "Comparing against your past uploads")
        similar_documents = []
        document_id = None
        try:
            similar_documents, document_id = await memory_service.remember_upload(
                user_id, result.filename, result.text
            )
        except Exception:
            pass

        return {
            "filename": result.filename,
            "text_content": result.text,
            "word_count": len(result.text.split()),
            "similar_documents": similar_documents,
            "document_id": document_id,
            "truncated": result.truncated,
        }

    except url_ingest.IngestError as e:
        # These messages are written for the user and say what to try next.
        raise HTTPException(status_code=422, detail=str(e))
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error ingesting URL: {str(e)}")
    finally:
        progress.clear(progress_id)


async def _weak_concepts_if_overlap(has_overlap: bool, user_id: str, text_content: str) -> Optional[List[str]]:
    """ROADMAP 3.2: when the upload overlapped prior material, fetch the
    user's weakest stored concepts that match this content, to pass to
    generation prompts as emphasis hints. Best-effort — never fails the
    generation itself."""
    if not has_overlap:
        return None
    try:
        return await memory_service.weak_concepts_for_text(user_id, text_content)
    except Exception:
        return None

@app.post("/generate-summary", response_model=SummaryResponse)
async def generate_summary(
    text_content: str = Form(...),
    summary_type: str = Form(...),  # "short", "bullet_points", "detailed"
    subject: Optional[str] = Form(None),
    progress_id: Optional[str] = Form(None),
    has_overlap: bool = Form(False),
    focus_concepts: Optional[str] = Form(None),  # JSON array of concept names
    # Free-text "what do you want to focus on" from the configure step. Passed
    # to the prompt verbatim — see ai_service._focus_note_block.
    focus_note: Optional[str] = Form(None),
    user_id: str = Depends(auth.get_current_user_id)
):
    """Generate summary from text content"""
    try:
        weak_concepts = await _weak_concepts_if_overlap(has_overlap, user_id, text_content)
        # Pretest-informed emphasis (ROADMAP_LEARNING 1): concepts the user
        # just missed on the pretest get extra coverage in the summary.
        if focus_concepts:
            try:
                parsed = json.loads(focus_concepts)
                if isinstance(parsed, list):
                    weak_concepts = list(dict.fromkeys((weak_concepts or []) + [str(c) for c in parsed]))
            except (json.JSONDecodeError, TypeError):
                pass
        summary = await ai_service.generate_summary(
            text_content=text_content,
            summary_type=summary_type,
            subject=subject,
            progress=progress.reporter(progress_id),
            weak_concepts=weak_concepts,
            focus_note=focus_note,
        )

        if summary_type == "bullet_points" and "concepts" in summary:
            summary_text = json.dumps({"concepts": summary["concepts"]})
            word_count = sum(
                len(c.get("title", "").split()) + len(c.get("explanation", "").split()) +
                sum(len(d.split()) for d in c.get("details", []))
                for c in summary["concepts"]
            )
        else:
            summary_text = summary["content"]
            word_count = len(summary["content"].split())

        return SummaryResponse(
            summary=summary_text,
            tags=summary.get("tags", []),
            summary_type=summary_type,
            word_count=word_count
        )
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error generating summary: {str(e)}")
    finally:
        progress.clear(progress_id)

@app.post("/generate-quiz", response_model=QuizResponse)
async def generate_quiz(
    text_content: str = Form(...),
    num_questions: int = Form(...),
    subject: str = Form(...),
    difficulty: str = Form(...),  # "easy", "medium", "hard"
    progress_id: Optional[str] = Form(None),
    has_overlap: bool = Form(False),
    focus_note: Optional[str] = Form(None),
    user_id: str = Depends(auth.get_current_user_id)
):
    """Generate quiz from text content"""
    try:
        weak_concepts = await _weak_concepts_if_overlap(has_overlap, user_id, text_content)
        quiz = await ai_service.generate_quiz(
            text_content=text_content,
            num_questions=num_questions,
            subject=subject,
            difficulty=difficulty,
            progress=progress.reporter(progress_id),
            weak_concepts=weak_concepts,
            focus_note=focus_note,
        )

        return QuizResponse(
            questions=quiz["questions"],
            total_questions=len(quiz["questions"]),
            difficulty=difficulty,
            subject=subject,
            estimated_time=len(quiz["questions"]) * 2  # 2 minutes per question
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error generating quiz: {str(e)}")
    finally:
        progress.clear(progress_id)

@app.post("/generate-flashcards", response_model=FlashcardResponse)
async def generate_flashcards(
    text_content: str = Form(...),
    num_cards: int = Form(...),
    subject: str = Form(...),
    card_type: str = Form(...),  # "definition", "concept", "fact", "mixed"
    document_id: Optional[str] = Form(None),
    progress_id: Optional[str] = Form(None),
    focus_note: Optional[str] = Form(None),
    user_id: str = Depends(auth.get_current_user_id)
):
    """Generate flashcards from text content"""
    try:
        flashcards = await ai_service.generate_flashcards(
            text_content=text_content,
            num_cards=num_cards,
            subject=subject,
            card_type=card_type,
            progress=progress.reporter(progress_id),
            focus_note=focus_note,
        )

        # Spaced repetition (ROADMAP 4.1): persist the set so the cards come
        # back for review at growing intervals. Best-effort — a DB failure
        # must never fail the generation the user is waiting on.
        try:
            db.save_flashcard_set(
                user_id, subject, card_type, flashcards["flashcards"],
                document_id=document_id,
            )
        except Exception:
            pass

        return FlashcardResponse(
            flashcards=flashcards["flashcards"],
            total_cards=len(flashcards["flashcards"]),
            subject=subject,
            card_type=card_type
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error generating flashcards: {str(e)}")
    finally:
        progress.clear(progress_id)

def _podcast_audio_url(podcast_id: str, audio_key: str, user_id: str) -> str:
    """Playback URL for an episode.

    Prefers a presigned S3 URL so the browser streams straight from S3 with
    range requests — scrubbing a long episode then costs the API nothing. When
    there's no bucket (local-disk storage) or presigning fails, falls back to
    the API route with a signed token in the query string.

    The token is necessary because a browser `<audio src>` cannot send an
    Authorization header, so the URL has to carry its own proof of access —
    exactly the job presigning does on the S3 path.
    """
    presigned = storage_service.presigned_url(audio_key)
    if presigned:
        return presigned

    # Absolute, because this URL goes into an <audio src> on the frontend
    # origin — a relative path would resolve against Next, not the API.
    base = os.getenv("PUBLIC_API_URL", "http://localhost:8000").rstrip("/")
    token = auth.make_media_token(podcast_id, user_id)
    query = urlencode({"token": token, "user_id": user_id})
    return f"{base}/podcasts/{podcast_id}/audio?{query}"

@app.post("/generate-podcast", response_model=PodcastResponse)
async def generate_podcast(
    text_content: str = Form(...),
    subject: str = Form(...),
    length: str = Form("medium"),  # "short", "medium", "long"
    document_id: Optional[str] = Form(None),
    progress_id: Optional[str] = Form(None),
    focus_note: Optional[str] = Form(None),
    user_id: str = Depends(auth.get_current_user_id)
):
    """Generate a two-speaker podcast episode from text content.

    The script and the audio fail independently on purpose. Writing and
    grounding the script is the expensive, valuable part; synthesis is a
    separate service that can be out of credit or misconfigured. So a
    synthesis failure still returns a 200 with the full script and an
    `audio_error` the player can show — the student gets a readable episode
    instead of losing a 30-second generation to someone else's billing.
    """
    if length not in ("short", "medium", "long"):
        raise HTTPException(status_code=400, detail="length must be short, medium, or long")

    reporter = progress.reporter(progress_id)

    try:
        script = await ai_service.generate_podcast_script(
            text_content=text_content,
            subject=subject,
            length=length,
            progress=reporter,
            focus_note=focus_note,
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error generating podcast: {str(e)}")
    finally:
        progress.clear(progress_id)

    segments = script["segments"]

    # The row is created before synthesis because its id names the audio
    # object in S3. A row whose audio_key stays null is the degraded episode.
    try:
        podcast_id = db.create_podcast(
            user_id, subject, script["title"], segments, document_id=document_id,
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error saving podcast: {str(e)}")

    audio_url: Optional[str] = None
    duration_seconds: Optional[int] = None
    audio_error: Optional[str] = None

    # Only the TTS key gates synthesis. Storage does not: without a bucket the
    # audio is written to local disk instead, so a machine with a Deepgram key
    # and no AWS setup still produces playable episodes.
    if not tts_service.is_configured():
        audio_error = (
            "Audio generation isn't set up on this server, so this episode is "
            "script-only."
        )
    else:
        try:
            audio, offsets, duration_seconds = await tts_service.synthesize_dialogue(
                segments, client=ai_service._client, progress=reporter,
            )
            key = await storage_service.put_bytes(
                storage_service.podcast_key(user_id, podcast_id), audio, "audio/mpeg",
            )
            # Fold the measured offsets back into the stored script, so a
            # later GET /podcasts/{id} serves the same exact timings rather
            # than the player having to re-estimate them.
            for segment, start in zip(segments, offsets):
                segment["start_seconds"] = start
            db.attach_podcast_audio(podcast_id, key, duration_seconds, segments)
            audio_url = _podcast_audio_url(podcast_id, key, user_id)
        except tts_service.TTSError as e:
            audio_error = e.user_message
        except storage_service.StorageError as e:
            audio_error = f"The episode was recorded but couldn't be stored: {e}"
        except Exception as e:
            audio_error = f"Audio synthesis failed: {str(e)}"
        finally:
            progress.clear(progress_id)

    return PodcastResponse(
        id=podcast_id,
        title=script["title"],
        subject=subject,
        segments=[PodcastSegment(**s) for s in segments],
        audio_url=audio_url,
        duration_seconds=duration_seconds,
        audio_error=audio_error,
        created_at=datetime.now(timezone.utc).isoformat().replace("+00:00", "Z"),
    )

@app.get("/me/podcasts", response_model=List[PodcastInfo])
async def get_my_podcasts(external_user_id: str = Depends(auth.get_current_user_id)):
    """All of the signed-in user's episodes, newest first"""
    try:
        return [PodcastInfo(**p) for p in db.list_podcasts(external_user_id)]
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error fetching podcasts: {str(e)}")

@app.get("/podcasts/{podcast_id}", response_model=PodcastResponse)
async def get_podcast(podcast_id: str, external_user_id: str = Depends(auth.get_current_user_id)):
    """One episode with its script and a fresh playback URL.

    The URL is minted per request rather than stored: presigned URLs expire,
    so a persisted one would be a link that works until it quietly doesn't.
    """
    podcast = db.get_podcast(podcast_id, external_user_id)
    if podcast is None:
        raise HTTPException(status_code=404, detail="Podcast not found")

    audio_key = podcast.get("audio_key")
    return PodcastResponse(
        id=podcast["id"],
        title=podcast["title"],
        subject=podcast["subject"],
        segments=[PodcastSegment(**s) for s in podcast["script"]],
        audio_url=(
            _podcast_audio_url(podcast["id"], audio_key, external_user_id)
            if audio_key else None
        ),
        duration_seconds=podcast.get("duration_seconds"),
        audio_error=None if audio_key else "This episode was never recorded as audio.",
        created_at=podcast["created_at"],
    )

@app.get("/podcasts/{podcast_id}/audio")
async def get_podcast_audio(
    podcast_id: str,
    token: Optional[str] = Query(None),
    user_id: Optional[str] = Query(None),
    authorization: Optional[str] = Header(None),
):
    """Stream an episode's audio through the API.

    The fallback for when presigning isn't available. Accepts either a normal
    bearer header or a signed `token` — a browser `<audio src>` can't set
    headers, so the token is what makes the element work at all. Ownership is
    still enforced either way: get_podcast is user-scoped, so a valid token for
    someone else's episode finds nothing.
    """
    if token and user_id:
        if not auth.verify_media_token(token, podcast_id, user_id):
            raise HTTPException(status_code=403, detail="Invalid or expired audio link")
        external_user_id = user_id
    else:
        external_user_id = await auth.get_current_user_id(authorization)

    podcast = db.get_podcast(podcast_id, external_user_id)
    if podcast is None or not podcast.get("audio_key"):
        raise HTTPException(status_code=404, detail="Podcast audio not found")

    try:
        audio = await storage_service.get_bytes(podcast["audio_key"])
    except storage_service.StorageError as e:
        raise HTTPException(status_code=502, detail=str(e))

    return Response(content=audio, media_type="audio/mpeg")

@app.get("/me/flashcards/due", response_model=DueFlashcardsResponse)
async def get_my_due_flashcards(external_user_id: str = Depends(auth.get_current_user_id)):
    """Cards due for review now (most overdue first) plus the total due
    count, for the review screen and the due-count badge"""
    try:
        return DueFlashcardsResponse(**db.get_due_flashcards(external_user_id))
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error fetching due flashcards: {str(e)}")

@app.get("/me/concepts/due", response_model=DueConceptReviewsResponse)
async def get_my_due_concepts(external_user_id: str = Depends(auth.get_current_user_id)):
    """Concepts whose spaced-review schedule says they're due (ROADMAP_LEARNING
    6), with their source document so the frontend can one-click a refresher
    tutor session on the stored content."""
    try:
        return DueConceptReviewsResponse(**db.get_due_concept_reviews(external_user_id))
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error fetching due concepts: {str(e)}")

@app.post("/flashcards/{card_id}/review", response_model=FlashcardReviewResponse)
async def review_flashcard(
    card_id: str,
    request: FlashcardReviewRequest,
    external_user_id: str = Depends(auth.get_current_user_id),
):
    """Apply one self-graded review ("again"/"hard"/"good"/"easy") to a card
    and return its new schedule"""
    if request.grade not in ("again", "hard", "good", "easy"):
        raise HTTPException(status_code=400, detail="grade must be one of: again, hard, good, easy")
    try:
        result = db.review_flashcard(card_id, external_user_id, request.grade)
        if result is None:
            raise HTTPException(status_code=404, detail="Flashcard not found")
        return FlashcardReviewResponse(**result)
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error recording review: {str(e)}")

# --- Pretesting (ROADMAP_LEARNING 1): retrieval before re-reading. A short
# --- quiz taken before the summary is shown; results calibrate the user's
# --- persistent concept mastery and flag weak spots in the summary.

@app.post("/pretest/start", response_model=PretestStartResponse)
async def pretest_start(request: PretestStartRequest, user_id: str = Depends(auth.get_current_user_id)):
    """Generate a short pretest (one multiple-choice question per extracted
    concept) to take before any summary is shown. Answers stay server-side."""
    try:
        if not request.text_content.strip():
            raise HTTPException(status_code=400, detail="text_content cannot be empty")
        result = await pretest_agent.start_pretest(
            user_id, request.text_content, request.subject, ai_service,
            document_id=request.document_id,
            progress=progress.reporter(request.progress_id),
        )
        return PretestStartResponse(**result)
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error starting pretest: {str(e)}")
    finally:
        progress.clear(request.progress_id)

@app.post("/pretest/submit", response_model=PretestSubmitResponse)
async def pretest_submit(request: PretestSubmitRequest, user_id: str = Depends(auth.get_current_user_id)):
    """Grade a pretest, write the per-concept outcomes into the user's
    persistent concept mastery (so a tutor session afterwards starts
    calibrated), and return the correction plus the missed concepts."""
    try:
        result = await pretest_agent.submit_pretest(request.pretest_id, user_id, request.answers)
        if result is None:
            raise HTTPException(status_code=404, detail="Pretest not found or expired")
        return PretestSubmitResponse(**result)
    except HTTPException:
        raise
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error submitting pretest: {str(e)}")

@app.post("/tutor/start", response_model=TutorStartResponse)
async def tutor_start(request: TutorStartRequest, user_id: str = Depends(auth.get_current_user_id)):
    """Start an adaptive tutor session: extract the concepts to teach from
    the uploaded material, initialize a per-concept knowledge state, and
    return the first question (without its answer — grading is server-side)."""
    try:
        # One session, one or many documents (ROADMAP_LEARNING 3): `documents`
        # is the general form, `text_content` the one-file shorthand. Collapse
        # them here so nothing downstream has to care which the client sent.
        if request.documents:
            sources = [
                {
                    "text_content": doc.text_content,
                    "filename": doc.filename,
                    "document_id": doc.document_id,
                }
                for doc in request.documents
                if doc.text_content and doc.text_content.strip()
            ]
        elif request.text_content.strip():
            sources = [{
                "text_content": request.text_content,
                "filename": request.subject,
                "document_id": request.document_id,
            }]
        else:
            sources = []
        if not sources:
            raise HTTPException(
                status_code=400,
                detail="provide documents[] or text_content with material to study",
            )
        if request.mode not in tutor_agent.MODES:
            raise HTTPException(status_code=400, detail="mode must be one of: " + ", ".join(tutor_agent.MODES))

        result = await tutor_agent.start_session(
            user_id, sources, request.subject, request.mode, ai_service,
            concepts_filter=request.concepts,
            progress=progress.reporter(request.progress_id),
        )
        return TutorStartResponse(**result)
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error starting tutor session: {str(e)}")
    finally:
        progress.clear(request.progress_id)

@app.post("/tutor/answer", response_model=TutorAnswerResponse)
async def tutor_answer(request: TutorAnswerRequest, user_id: str = Depends(auth.get_current_user_id)):
    """Submit an answer to the current tutor question. Grades it, diagnoses
    wrong answers, updates the knowledge state, and returns either the next
    question (targeting the weakest concept at a calibrated difficulty) or
    the session summary."""
    try:
        result = await tutor_agent.submit_answer(
            request.session_id, user_id, request.answer, ai_service,
            confidence=request.confidence,
        )
        if result is None:
            raise HTTPException(status_code=404, detail="Tutor session not found or expired")
        return TutorAnswerResponse(**result)
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing answer: {str(e)}")

@app.post("/tutor/wrap", response_model=TutorWrapResponse)
async def tutor_wrap(request: TutorWrapRequest, user_id: str = Depends(auth.get_current_user_id)):
    """End an active tutor session early at the student's request (the
    soft-checkpoint "wrap up" action) and return its summary."""
    try:
        summary = tutor_agent.wrap_session(request.session_id, user_id)
        if summary is None:
            raise HTTPException(status_code=404, detail="Tutor session not found or expired")
        return TutorWrapResponse(summary=summary)
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error wrapping tutor session: {str(e)}")

@app.get("/tutor/session/{session_id}", response_model=TutorStartResponse)
async def tutor_get_session(session_id: str, user_id: str = Depends(auth.get_current_user_id)):
    """Current state of an active tutor session, for resuming the UI after a
    page refresh: the pending question (answer stays server-side) and the
    per-concept knowledge state."""
    state = tutor_agent.get_session_state(session_id, user_id)
    if state is None:
        raise HTTPException(status_code=404, detail="Tutor session not found or already finished")
    return TutorStartResponse(**state)

# --- Voice roleplay (ROADMAP_HONEN 4) ---
#
# The session is created over plain HTTP first, then a websocket attaches to
# it. Scenario generation is a 10-30s LLM pipeline that wants the existing
# progress-polling UX, and doing it inside the socket handshake would mean a
# half-minute of silence on a connection the client can't show progress for.

# How long the client has to send its auth frame before the socket is closed.
# Short on purpose: an unauthenticated socket costs a connection slot, and a
# real client sends this frame immediately.
WS_AUTH_TIMEOUT_SECONDS = 5

# Close codes. 4401/4404 mirror HTTP 401/404 in the application range, since
# the WS close frame is the only channel available once the socket is open.
WS_CLOSE_UNAUTHORIZED = 4401
WS_CLOSE_NOT_FOUND = 4404


@app.post("/roleplay/start", response_model=RoleplayStartResponse)
async def roleplay_start(request: RoleplayStartRequest, user_id: str = Depends(auth.get_current_user_id)):
    """Generate a grounded scene and open a roleplay session.

    Refuses at the door when STT is unconfigured. This is deliberately NOT the
    /generate-podcast independent-failure pattern: TTS degrades to captions and
    the scene still works, but STT is the *input* channel — without it there is
    no scene, only a character talking to nobody. Checking before the LLM runs
    means a missing key costs nothing rather than surfacing after 30 seconds of
    scenario generation.
    """
    if not stt_service.is_configured():
        raise HTTPException(
            status_code=503,
            detail="Voice roleplay isn't set up on this server (DEEPGRAM_API_KEY is missing).",
        )

    try:
        if request.documents:
            sources = [
                {
                    "text_content": doc.text_content,
                    "filename": doc.filename,
                    "document_id": doc.document_id,
                }
                for doc in request.documents
                if doc.text_content and doc.text_content.strip()
            ]
        elif request.text_content.strip():
            sources = [{
                "text_content": request.text_content,
                "filename": request.subject,
                "document_id": request.document_id,
            }]
        else:
            sources = []
        if not sources:
            raise HTTPException(
                status_code=400,
                detail="provide documents[] or text_content with material to study",
            )

        result = await roleplay_agent.start_session(
            user_id, sources, request.subject, ai_service,
            concept=request.concept,
            progress=progress.reporter(request.progress_id),
        )
        return RoleplayStartResponse(**result)
    except HTTPException:
        raise
    except Exception as e:
        # Fatal, not degraded: a canned fallback scene would be ungrounded,
        # and an ungrounded scene defeats the premise of the feature.
        raise HTTPException(status_code=500, detail=f"Error starting roleplay session: {str(e)}")
    finally:
        progress.clear(request.progress_id)


@app.post("/roleplay/end", response_model=RoleplayResultResponse)
async def roleplay_end(request: RoleplayEndRequest, user_id: str = Depends(auth.get_current_user_id)):
    """Grade the scene and complete the session.

    Over plain HTTP rather than the socket so a dead connection never costs the
    student their result — this is the path the client falls back to when the
    websocket drops mid-scene.
    """
    session = roleplay_agent._load_session(request.session_id, user_id)
    if session is None:
        raise HTTPException(status_code=404, detail="Roleplay session not found or already finished")

    try:
        result = await roleplay_agent.grade_session(request.session_id, session, ai_service)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error grading roleplay session: {str(e)}")
    return RoleplayResultResponse(**result)


@app.get("/roleplay/{session_id}/result", response_model=RoleplayResultResponse)
async def roleplay_result(session_id: str, user_id: str = Depends(auth.get_current_user_id)):
    """A finished scene's rubric result and transcript."""
    result = await asyncio.to_thread(db.get_roleplay_result, session_id, user_id)
    if result is None:
        raise HTTPException(status_code=404, detail="Roleplay result not found")
    return RoleplayResultResponse(**result)


async def _ws_send(websocket: WebSocket, payload: dict) -> bool:
    """Send one JSON control frame. False if the client is gone.

    Every send goes through here: with no test suite, a client closing
    mid-turn is the likeliest source of log noise, and an unguarded send in
    the middle of a turn would take down the whole handler.
    """
    try:
        await websocket.send_json(payload)
        return True
    except (WebSocketDisconnect, RuntimeError):
        return False


# Binary frame size for the MP3 down-channel. Small enough that a barge-in
# lands within a frame or two, large enough that a 100KB clip isn't hundreds of
# sends.
WS_AUDIO_CHUNK_BYTES = 16 * 1024


async def _stream_turn_audio(websocket: WebSocket, live: dict, turn: dict) -> None:
    """Synthesize one reply and ship it as binary frames, then audio_end.

    Degrades per-channel, following /generate-podcast: a synthesis failure
    still sends `audio_end` so the client's turn state advances and the line
    stands as text. `notice{degraded:true}` is sent **once** per session — a
    broken key would otherwise add a stall and a toast to every single turn.
    """
    turn_id = turn["turn_id"]

    async def _finish(degraded_message: Optional[str] = None) -> None:
        if degraded_message and not live["tts_degraded"]:
            live["tts_degraded"] = True
            await _ws_send(websocket, {
                "type": "notice", "code": "tts_degraded",
                "message": degraded_message, "degraded": True,
            })
        # Always sent, even with zero audio frames: the client keys "the
        # character has finished speaking" off this frame, so skipping it on
        # the degraded path would strand the UI mid-turn forever.
        await _ws_send(websocket, {"type": "audio_end", "turn_id": turn_id})

    if live["tts_degraded"] or not tts_service.is_configured():
        await _finish(
            None if live["tts_degraded"] else
            "Audio isn't set up on this server, so the character's lines will "
            "appear as text."
        )
        return

    task = asyncio.create_task(
        tts_service.synthesize_turn(turn["reply"], client=ai_service._client)
    )
    live["tts_task"] = task

    try:
        audio = await task
    except asyncio.CancelledError:
        # Barge-in: the student started talking over this line. Abandon it
        # quietly — the client already dropped the buffer.
        await _ws_send(websocket, {"type": "audio_end", "turn_id": turn_id})
        return
    except tts_service.TTSError as e:
        await _finish(e.user_message)
        return
    except Exception as e:
        await _finish(f"Audio synthesis failed: {e}")
        return
    finally:
        live["tts_task"] = None

    for start in range(0, len(audio), WS_AUDIO_CHUNK_BYTES):
        # A turn that's been superseded by a barge-in stops shipping frames.
        if live["current_turn_id"] != turn_id:
            break
        try:
            await websocket.send_bytes(audio[start:start + WS_AUDIO_CHUNK_BYTES])
        except (WebSocketDisconnect, RuntimeError):
            return

    await _finish()


@app.websocket("/roleplay/live/{session_id}")
async def roleplay_live(websocket: WebSocket, session_id: str):
    """The live roleplay channel — the first websocket in this codebase.

    Framing rule: binary frames are always audio, text frames are always JSON
    control. Direction disambiguates the two binary formats (PCM up, MP3 down),
    so nothing is base64'd.

    Auth is a first-frame protocol rather than a query parameter: a token in
    the URL lands in uvicorn's access logs and the browser's history, and a WS
    URL is not treated as a secret by anything that handles it.
    """
    await websocket.accept()

    # Authentication and ownership are separate checks with separate codes.
    # A foreign session is 4404, not 4403 — indistinguishable from one that
    # doesn't exist, matching the rule the HTTP loaders follow.
    try:
        opening = await asyncio.wait_for(
            websocket.receive_json(), timeout=WS_AUTH_TIMEOUT_SECONDS,
        )
    except (asyncio.TimeoutError, WebSocketDisconnect, ValueError):
        await websocket.close(code=WS_CLOSE_UNAUTHORIZED)
        return

    if not isinstance(opening, dict) or opening.get("type") != "auth":
        await websocket.close(code=WS_CLOSE_UNAUTHORIZED)
        return

    user_id = await auth.verify_bearer_token(str(opening.get("token") or ""))
    if not user_id:
        await websocket.close(code=WS_CLOSE_UNAUTHORIZED)
        return

    session = roleplay_agent._load_session(session_id, user_id)
    if session is None:
        await websocket.close(code=WS_CLOSE_NOT_FOUND)
        return

    live = session["live"]
    live["websocket"] = websocket
    persisted = session["persisted"]

    # current_turn_id resumes from the transcript rather than restarting at 0:
    # a reconnect mid-scene must not reuse turn ids the client has already
    # seen, or its "discard frames from a stale turn" rule discards live ones.
    live["current_turn_id"] = max(
        (turn.get("turn_id", 0) for turn in persisted["transcript"]), default=0,
    )

    await _ws_send(websocket, {
        "type": "ready",
        "scenario": roleplay_agent.public_scenario(persisted["scenario"]),
        "transcript": persisted["transcript"],
        "turns_taken": persisted["turns_taken"],
    })

    # The scene's turn queue. The socket reader and the Flux pump both feed it,
    # so a spoken turn and a typed one land on exactly the same code path
    # below — the transport is the only thing that differs.
    utterances: asyncio.Queue = asyncio.Queue()
    flux_task: Optional[asyncio.Task] = None

    async def _pump_flux() -> None:
        """Translate Flux turn events onto the wire protocol.

        Flux carries end-of-turn detection itself, so this is a mapping, not a
        turn-taking algorithm: EndOfTurn arrives with the final transcript
        already attached, which is why there is no separate final-pass request.
        """
        flux = live.get("flux")
        if flux is None:
            return
        try:
            async for event in flux.events():
                name = event.get("event")
                turn_index = event.get("turn_index", 0)

                if name == "Error":
                    if not live["stt_degraded"]:
                        live["stt_degraded"] = True
                        await _ws_send(websocket, {
                            "type": "notice", "code": "stt_unavailable",
                            "message": (
                                "We lost the speech service. You can keep going "
                                "by typing instead."
                            ),
                            "degraded": True,
                        })
                    continue

                if name in ("Update", "TurnResumed"):
                    # Advisory. `transcript` is the whole turn so far, not a
                    # delta, which is why the client replaces rather than
                    # appends. TurnResumed means Flux withdrew a draft
                    # end-of-turn, so the same non-final frame reverts it.
                    await _ws_send(websocket, {
                        "type": "transcript", "text": event.get("transcript", ""),
                        "final": False, "turn_id": turn_index,
                    })
                elif name == "EndOfTurn":
                    text = (event.get("transcript") or "").strip()
                    await _ws_send(websocket, {
                        "type": "transcript", "text": text,
                        "final": True, "turn_id": turn_index,
                    })
                    # An empty EndOfTurn is silence Flux mistook for speech.
                    # Running a turn on "" produces a confused reply to
                    # nothing, so drop it rather than invent a transcript.
                    if text:
                        await utterances.put(text)
        except asyncio.CancelledError:
            raise
        except Exception:
            logger.debug("roleplay: flux pump ended for %s", session_id, exc_info=True)

    async def _read_socket() -> None:
        """Route inbound frames: binary is audio, text is JSON control."""
        nonlocal flux_task
        while True:
            try:
                raw = await websocket.receive()
            except (WebSocketDisconnect, RuntimeError):
                await utterances.put(None)
                return

            if raw.get("type") == "websocket.disconnect":
                await utterances.put(None)
                return

            # Binary: 80ms of 16kHz mono PCM, straight upstream.
            audio = raw.get("bytes")
            if audio is not None:
                # Server-side gate as well as the client's. The client already
                # drops frames while muted, but a stale or buggy client would
                # otherwise bill us for uploaded silence — and Flux charges for
                # streamed audio whether anyone is talking or not.
                if live["stt_degraded"] or not live["mic_open"]:
                    continue
                flux = live.get("flux")
                if flux is None:
                    flux = stt_service.FluxSession(
                        keyterms=(persisted["scenario"] or {}).get("grounding_concepts")
                    )
                    live["flux"] = flux
                    flux_task = asyncio.create_task(_pump_flux())
                try:
                    await flux.send_audio(audio)
                except stt_service.STTError as e:
                    if not live["stt_degraded"]:
                        live["stt_degraded"] = True
                        await _ws_send(websocket, {
                            "type": "notice", "code": e.code,
                            "message": e.user_message, "degraded": True,
                        })
                continue

            text_frame = raw.get("text")
            if text_frame is None:
                continue
            try:
                message = json.loads(text_frame)
            except (ValueError, TypeError):
                continue
            if not isinstance(message, dict):
                continue

            kind = message.get("type")

            if kind == "end_session":
                await utterances.put(None)
                return

            if kind == "barge_in":
                # Bump the turn id so any in-flight frame loop stops, and
                # cancel the synthesis task. The client has already called
                # source.stop() and dropped its buffer — because a turn is
                # decoded only at audio_end, there's no decoder state to
                # unwind, just a byte buffer to discard.
                live["current_turn_id"] += 1
                task = live.get("tts_task")
                if task is not None and not task.done():
                    task.cancel()
                continue

            if kind == "utterance":
                text = str(message.get("text") or "").strip()
                if text:
                    await utterances.put(text)
                continue

            # mic_open / mic_close are "start/stop sending frames upstream" —
            # a mute button and the push-to-talk gate, NOT utterance
            # boundaries. Flux owns endpointing and keeps its socket across
            # both; only end_session sends CloseStream. Closing the mic
            # deliberately does not close the upstream socket: that would
            # discard the conversational state Flux's end-of-turn model
            # carries, which is the whole reason the socket is per-session.
            if kind in ("mic_open", "mic_close"):
                live["mic_open"] = kind == "mic_open"

    reader_task = asyncio.create_task(_read_socket())

    try:
        while True:
            text = await utterances.get()
            if text is None:
                break

            if not await _ws_send(websocket, {
                "type": "thinking", "turn_id": live["current_turn_id"] + 1,
            }):
                break

            try:
                turn = await roleplay_agent.handle_utterance(
                    session_id, session, text, ai_service,
                )
            except Exception as e:
                logger.warning("roleplay: turn failed for %s", session_id, exc_info=True)
                await _ws_send(websocket, {
                    "type": "notice", "code": "turn_failed",
                    "message": "That turn didn't go through. Try saying it again.",
                    "degraded": False,
                })
                continue

            # reply_text must precede its audio, and audio_end must follow it.
            # The client keys its turn state off that ordering.
            if not await _ws_send(websocket, {
                "type": "reply_text", "text": turn["reply"], "turn_id": turn["turn_id"],
            }):
                break

            await _stream_turn_audio(websocket, live, turn)

            if turn["nudge"]:
                await _ws_send(websocket, {
                    "type": "notice", "code": "soft_nudge",
                    "message": "You've covered a lot here — wrap up whenever you're ready.",
                    "degraded": False,
                })

            if turn["done"]:
                break

        # Grading runs on the way out so a client that closes cleanly still
        # gets its result over the socket. A client that vanishes instead
        # reaches the same result through POST /roleplay/end.
        result = await roleplay_agent.grade_session(session_id, session, ai_service)
        await _ws_send(websocket, {"type": "graded", "result": result})
    except Exception:
        logger.warning("roleplay: socket handler failed for %s", session_id, exc_info=True)
        await _ws_send(websocket, {
            "type": "error", "code": "internal",
            "message": "Something went wrong with this scene.",
        })
    finally:
        live["websocket"] = None

        # Close the upstream Flux socket before anything else: it is the one
        # resource here that costs money while it stays open.
        flux = live.get("flux")
        live["flux"] = None
        if flux is not None:
            try:
                await flux.close()
            except Exception:
                pass

        for task in (reader_task, flux_task, live.get("tts_task")):
            if task is not None and not task.done():
                task.cancel()
        live["tts_task"] = None

        try:
            await websocket.close()
        except Exception:
            pass


# --- Documents library (ROADMAP 3.1): the memory layer's stored uploads,
# --- made user-visible so material can be re-studied without re-uploading.

@app.get("/me/documents", response_model=List[DocumentInfo])
async def get_my_documents(external_user_id: str = Depends(auth.get_current_user_id)):
    """All of the signed-in user's stored uploads, newest first"""
    try:
        documents = db.list_documents(external_user_id)
        return [DocumentInfo(**d) for d in documents]
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error fetching documents: {str(e)}")

@app.get("/documents/{document_id}/content", response_model=DocumentContent)
async def get_document_content(document_id: str, external_user_id: str = Depends(auth.get_current_user_id)):
    """Reassembled text of a stored upload, for studying it again"""
    content = db.get_document_content(document_id, external_user_id)
    if content is None:
        raise HTTPException(status_code=404, detail="Document not found")
    return DocumentContent(**content)

# --- Viewing the original upload -------------------------------------------
# The extracted text is what gets studied, but extraction is lossy in ways the
# student can't see. These routes serve the file they actually uploaded.


async def _document_source_or_404(document_id: str, external_user_id: str, require_pdf: bool = False):
    """Ownership-scoped lookup of a document's stored file.

    Raises 404 for a foreign document, an unknown one, or one with no stored
    file — all indistinguishable on purpose, so the routes never confirm that
    someone else's document exists.
    """
    document = db.get_document_source(document_id, external_user_id)
    if document is None or not document.get("source_key"):
        raise HTTPException(status_code=404, detail="No original file for this document")
    if require_pdf and document.get("source_content_type") != db.PDF_CONTENT_TYPE:
        raise HTTPException(status_code=404, detail="This document is not a PDF")
    return document


async def _resolve_media_user(
    resource_id: str, token: Optional[str], user_id: Optional[str],
    authorization: Optional[str],
) -> str:
    """Authenticate a binary-media request by signed token or bearer header.

    An <img src> / <a download> cannot set an Authorization header, so the
    token carries the grant in the URL instead — the same constraint that
    drives podcast audio playback. Authentication only; ownership is a
    separate check the caller must still make.
    """
    if token and user_id:
        if not auth.verify_media_token(token, resource_id, user_id):
            raise HTTPException(status_code=403, detail="Invalid or expired link")
        return user_id
    return await auth.get_current_user_id(authorization)


def _document_media_url(path: str, document_id: str, user_id: str, token: str) -> str:
    """Absolute, self-authenticating URL for one of the media routes below.

    Absolute because it lands in an <img src> on the frontend origin, where a
    relative path would resolve against Next rather than the API.
    """
    base = os.getenv("PUBLIC_API_URL", "http://localhost:8000").rstrip("/")
    query = urlencode({"token": token, "user_id": user_id})
    return f"{base}/documents/{document_id}/{path}?{query}"


@app.get("/documents/{document_id}/original/meta", response_model=DocumentOriginalMeta)
async def get_document_original_meta(
    document_id: str, external_user_id: str = Depends(auth.get_current_user_id),
):
    """What the viewer needs before rendering: is there a file, can it be
    paged through, how many pages, and a token for the media URLs.

    Every failure collapses to `available: false` rather than an error status.
    A missing or unreadable original is a normal state — documents uploaded
    before this feature and documents ingested from a URL have no file at all
    — so the UI shows one honest line instead of an error.
    """
    document = db.get_document_source(document_id, external_user_id)
    if document is None:
        raise HTTPException(status_code=404, detail="Document not found")

    filename = document["filename"]
    content_type = document.get("source_content_type")
    is_pdf = content_type == db.PDF_CONTENT_TYPE

    if not document.get("source_key"):
        return DocumentOriginalMeta(available=False, is_pdf=False, filename=filename)

    pages = None
    if is_pdf:
        # The one place the PDF is opened to count pages. A corrupt or
        # password-protected file fails here; it stays downloadable, because
        # unreadable to us doesn't mean unreadable to the student.
        try:
            data = await storage_service.get_bytes(document["source_key"])
            pages = await asyncio.to_thread(pdf_render.page_count, data)
        except (storage_service.StorageError, pdf_render.RenderError):
            logger.warning("Could not read original for document %s", document_id, exc_info=True)
            is_pdf = False

    # One token for the whole viewing session: the client substitutes {page}
    # and every image load reuses the same grant.
    token = auth.make_media_token(document_id, external_user_id)
    return DocumentOriginalMeta(
        available=True,
        is_pdf=is_pdf,
        filename=filename,
        page_count=pages,
        content_type=content_type,
        page_url_template=(
            _document_media_url("page/{page}", document_id, external_user_id, token)
            if is_pdf else None
        ),
        download_url=_document_media_url("original", document_id, external_user_id, token),
    )


@app.get("/documents/{document_id}/page/{page_number}")
async def get_document_page(
    document_id: str,
    page_number: int,
    token: Optional[str] = Query(None),
    user_id: Optional[str] = Query(None),
    authorization: Optional[str] = Header(None),
):
    """One page of a stored PDF, rendered to PNG on demand."""
    external_user_id = await _resolve_media_user(document_id, token, user_id, authorization)
    document = await _document_source_or_404(document_id, external_user_id, require_pdf=True)

    try:
        data = await storage_service.get_bytes(document["source_key"])
        png = await asyncio.to_thread(pdf_render.render_page, data, page_number)
    except pdf_render.RenderError as e:
        raise HTTPException(status_code=404, detail=str(e))
    except storage_service.StorageError as e:
        raise HTTPException(status_code=502, detail=str(e))

    # Caching is load-bearing, not polish: without it every page flip re-fetches
    # the whole PDF from S3 and re-renders it. max-age matches the media token's
    # lifetime, and `private` keeps a shared cache from ever holding it.
    return Response(
        content=png,
        media_type="image/png",
        headers={"Cache-Control": f"private, max-age={auth.MEDIA_TOKEN_TTL_SECONDS}"},
    )


@app.get("/documents/{document_id}/original")
async def get_document_original(
    document_id: str,
    token: Optional[str] = Query(None),
    user_id: Optional[str] = Query(None),
    authorization: Optional[str] = Header(None),
):
    """Download the file as uploaded.

    The only way to see a DOCX or PPTX — rendering those would need LibreOffice
    — and a useful escape hatch for PDFs whose pages we can't render.
    """
    external_user_id = await _resolve_media_user(document_id, token, user_id, authorization)
    document = await _document_source_or_404(document_id, external_user_id)

    # Unlike a rendered page, this is a direct object fetch, so S3 can serve it
    # itself and 25 MB never crosses this process.
    presigned = storage_service.presigned_url(document["source_key"])
    if presigned:
        return RedirectResponse(presigned)

    try:
        data = await storage_service.get_bytes(document["source_key"])
    except storage_service.StorageError as e:
        raise HTTPException(status_code=502, detail=str(e))

    # The filename is client-supplied and has not been re-emitted anywhere
    # since upload; quote it so it can't inject header syntax.
    safe_name = os.path.basename(document["filename"]).replace('"', "")
    return Response(
        content=data,
        media_type=document.get("source_content_type") or "application/octet-stream",
        headers={"Content-Disposition": f'attachment; filename="{safe_name}"'},
    )


@app.delete("/documents/{document_id}")
async def delete_document(document_id: str, external_user_id: str = Depends(auth.get_current_user_id)):
    """Delete a stored upload and its chunks (ownership-scoped)"""
    try:
        deleted = db.delete_document(document_id, external_user_id)
        if not deleted:
            raise HTTPException(status_code=404, detail="Document not found")
        return {"deleted": True}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error deleting document: {str(e)}")

@app.post("/subjects", response_model=Subject)
async def create_subject(request: CreateSubjectRequest, external_user_id: str = Depends(auth.get_current_user_id)):
    """Create a subject/project for the signed-in user, or return the
    existing one if a subject with this name already exists"""
    try:
        if not request.name.strip():
            raise HTTPException(status_code=400, detail="Subject name cannot be empty")
        subject = db.create_subject(external_user_id, request.name)
        return Subject(**subject)
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error creating subject: {str(e)}")

@app.get("/subjects", response_model=List[Subject])
async def get_subjects(external_user_id: str = Depends(auth.get_current_user_id)):
    """List all subjects owned by the signed-in user"""
    try:
        subjects = db.list_subjects(external_user_id)
        return [Subject(**s) for s in subjects]
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error fetching subjects: {str(e)}")

@app.delete("/subjects/{subject_id}")
async def delete_subject(subject_id: str, external_user_id: str = Depends(auth.get_current_user_id)):
    """Delete a subject owned by the signed-in user. Past attempts that
    referenced it survive and fall into 'Uncategorized' in subject-grouped
    views (ON DELETE SET NULL on quiz_attempts.subject_id)."""
    try:
        deleted = db.delete_subject(subject_id, external_user_id)
        if not deleted:
            raise HTTPException(status_code=404, detail="Subject not found")
        return {"deleted": True}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error deleting subject: {str(e)}")

@app.post("/check-answers", response_model=AnswerCheckResponse)
async def check_answers(request: AnswerCheckRequest, external_user_id: str = Depends(auth.get_current_user_id)):
    """Check user answers, score the quiz, and persist the attempt"""
    try:
        if len(request.questions) != len(request.user_answers):
            raise HTTPException(status_code=400, detail="Answer count mismatch")

        user_id = db.get_or_create_user(external_user_id)

        result = db.record_quiz_attempt(
            subject_id=request.subject_id,
            difficulty=request.difficulty,
            questions=request.questions,
            user_answers=request.user_answers,
            user_id=user_id,
        )

        return AnswerCheckResponse(**result)

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error checking answers: {str(e)}")

@app.get("/quiz-attempts/{attempt_id}/breakdown", response_model=AttemptBreakdownResponse)
async def get_attempt_breakdown(attempt_id: str, user_id: str = Depends(auth.get_current_user_id)):
    """Real per-category and per-difficulty performance for a single completed attempt"""
    try:
        breakdown = db.get_attempt_breakdown(attempt_id)
        return AttemptBreakdownResponse(**breakdown)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error fetching breakdown: {str(e)}")

@app.get("/me/stats", response_model=UserStatsResponse)
async def get_my_stats(external_user_id: str = Depends(auth.get_current_user_id)):
    """Aggregate quiz-history stats for the signed-in user, for the profile screen"""
    try:
        stats = db.get_user_stats(external_user_id)
        return UserStatsResponse(**stats)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error fetching stats: {str(e)}")

@app.get("/me/analytics", response_model=UserAnalyticsResponse)
async def get_my_analytics(external_user_id: str = Depends(auth.get_current_user_id)):
    """Chart-ready datasets for the signed-in user: score trend, accuracy
    by category/difficulty, and quiz distribution by subject"""
    try:
        analytics = db.get_user_analytics(external_user_id)
        return UserAnalyticsResponse(**analytics)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error fetching analytics: {str(e)}")

@app.get("/me/recent-attempts", response_model=List[RecentAttempt])
async def get_my_recent_attempts(
    limit: int = Query(20, ge=1, le=200),
    external_user_id: str = Depends(auth.get_current_user_id),
):
    """Lightweight past-attempts list. The default is a short preview; the
    scores page asks for more to show a full history."""
    try:
        attempts = db.get_recent_attempts(external_user_id, limit=limit)
        return [RecentAttempt(**a) for a in attempts]
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error fetching recent attempts: {str(e)}")

@app.get("/quiz-attempts/{attempt_id}/recap", response_model=AttemptRecapResponse)
async def get_attempt_recap(attempt_id: str, external_user_id: str = Depends(auth.get_current_user_id)):
    """Full read-only recap of a past attempt, scoped to the requesting user"""
    recap = db.get_attempt_recap(attempt_id, external_user_id)
    if recap is None:
        raise HTTPException(status_code=404, detail="Attempt not found")
    return AttemptRecapResponse(**recap)

def extract_text_from_docx(file_path: str) -> str:
    """Extract text content from a DOCX file"""
    try:
        document = docx.Document(file_path)
        parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text]

        for table in document.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text:
                        parts.append(cell.text)

        return "\n".join(parts).strip()

    except Exception as e:
        raise Exception(f"Error extracting text from DOCX: {str(e)}")

def extract_text_from_pptx(file_path: str) -> str:
    """Extract text content from a PPTX file"""
    try:
        presentation = Presentation(file_path)
        parts = []

        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text:
                    parts.append(shape.text_frame.text)

        return "\n".join(parts).strip()

    except Exception as e:
        raise Exception(f"Error extracting text from PPTX: {str(e)}")

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000) 